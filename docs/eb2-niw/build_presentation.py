#!/usr/bin/env python3
"""Generates the Genesis AI pitch presentation — white theme."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (white theme) ─────────────────────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white
BLACK   = RGBColor(0x11, 0x11, 0x11)   # near-black text
YELLOW  = RGBColor(0xF5, 0xE6, 0x42)   # Genesis yellow
DGRAY   = RGBColor(0x44, 0x44, 0x44)   # body text
MGRAY   = RGBColor(0x88, 0x88, 0x88)   # muted text
LGRAY   = RGBColor(0xF0, 0xF0, 0xF0)   # box fill
BORDER  = RGBColor(0xDD, 0xDD, 0xDD)   # box border
INK     = RGBColor(0x1A, 0x1A, 0x1A)   # dark box bg
INKTEXT = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def tb(slide, text, l, t, w, h, size,
       bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def multiline_tb(slide, lines, l, t, w, h, default_size,
                 default_color=BLACK, default_bold=False, align=PP_ALIGN.LEFT):
    """lines = list of (text, size, bold, color) tuples"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill=LGRAY, line=BORDER, lw=1.0):
    from pptx.util import Pt as Ptx
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Ptx(lw)
    else:
        shp.line.fill.background()
    return shp


def accent_bar(slide, l, t, h, color=YELLOW, w=0.06):
    rect(slide, l, t, w, h, fill=color, line=None)


def section_heading(slide, text, l=0.55, t=0.35):
    accent_bar(slide, l, t, 0.65)
    tb(slide, text, l + 0.15, t, 9.0, 0.75, 40, bold=True, color=BLACK)


def bullet_list(slide, items, l, t, w, h, size, color=DGRAY,
                header=None, header_color=BLACK, header_size=None, gap=0.05):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]
        first = False
        r = p.add_run()
        r.text = header
        r.font.size = Pt(header_size or size + 1)
        r.font.bold = True
        r.font.color.rgb = header_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"  •  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
# Yellow accent block left
rect(s, 0.0, 0.0, 0.18, 7.5, fill=YELLOW, line=None)
tb(s, "Genesis", 0.55, 1.5, 9.0, 1.8, 100, bold=True, color=BLACK)
tb(s, "⚡", 9.15, 1.45, 1.0, 1.0, 52, color=YELLOW)
tb(s, "INTELLIGENT BRAND DESIGN", 0.55, 3.35, 10.0, 0.6, 20,
   bold=True, color=MGRAY)
tb(s, "The world's first AI creative director for growing businesses.",
   0.55, 4.1, 10.0, 0.7, 24, color=DGRAY)
tb(s, "Brand creation  ·  Ongoing strategy  ·  Assets on demand  ·  Physical products",
   0.55, 5.0, 11.0, 0.5, 15, color=MGRAY, italic=True)
tb(s, "Darrell Irwin & Saulo Oliveira  ·  genesisbrands.ai",
   0.55, 6.5, 8.0, 0.5, 13, color=MGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Problem
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Problem")
tb(s, "Brand is the most important part of any start-up's journey.",
   0.55, 1.35, 10.0, 0.65, 28, bold=True, color=BLACK)
tb(s, "Yet it remains the hardest to execute — constrained by budget, time, and confidence.",
   0.55, 2.1, 9.5, 0.6, 22, color=DGRAY)

stats = [
    ("5.5M", "new US businesses registered every year"),
    ("<10%", "can afford professional branding in year one"),
    ("$15K–$75K+", "what a traditional agency charges for a brand"),
    ("Weeks–months", "how long agencies take to deliver"),
]
for i, (num, label) in enumerate(stats):
    x = 0.55 + (i % 2) * 6.3
    y = 3.0 + (i // 2) * 1.7
    r = rect(s, x, y, 5.9, 1.45, fill=LGRAY, line=BORDER)
    tb(s, num,   x + 0.2, y + 0.1, 5.5, 0.7, 30, bold=True, color=BLACK)
    tb(s, label, x + 0.2, y + 0.82, 5.5, 0.5, 15, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Solution
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Solution")
tb(s, "A trained AI creative director — not a template generator.",
   0.55, 1.35, 11.5, 0.65, 26, bold=True, color=BLACK)
tb(s, "Founders answer a structured discovery questionnaire. Genesis AI — the Brain Engine —\nconducts intelligent discovery, derives brand strategy, and orchestrates specialist agents\nto produce a complete brand in minutes.",
   0.55, 2.1, 11.5, 1.5, 19, color=DGRAY)

boxes = [
    ("Brand Creation",  "Complete brand identity in minutes:\nlogo, copy, colour, typography,\nbrand book — three directions."),
    ("Ongoing Strategy","Subscription access to the AI\ncreative director that knows\nyour brand inside out."),
    ("Assets on Demand","Generate social media posts,\ncarousels, graphics and more\nwhenever you need them."),
    ("Physical Products","Order branded business cards,\nbanners, merch — printed and\ndelivered via our partners."),
]
for i, (title, body) in enumerate(boxes):
    x = 0.4 + i * 3.25
    r = rect(s, x, 3.9, 3.05, 2.95, fill=LGRAY, line=BORDER)
    accent_bar(s, x, 3.9, 0.35, color=YELLOW)
    tb(s, title, x + 0.2, 4.0, 2.7, 0.45, 15, bold=True, color=BLACK)
    tb(s, body,  x + 0.2, 4.55, 2.7, 2.2, 14, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — The Brain Engine (core differential)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "The Differential")
tb(s, "The Brain Engine — a trained second brain, not a prompt wrapper.",
   0.55, 1.35, 12.0, 0.6, 24, bold=True, color=BLACK)

# Left column — what it is NOT
r1 = rect(s, 0.5, 2.1, 5.7, 4.8, fill=LGRAY, line=BORDER)
tb(s, "What everyone else does", 0.7, 2.25, 5.3, 0.5, 16, bold=True, color=MGRAY)
for i, line in enumerate([
    "Ask a few style questions",
    "Match to a template library",
    "Generate one output",
    "Client iterates alone",
    "No brand memory",
    "No quality inspection",
    "Transaction ends at delivery",
]):
    tb(s, f"✗  {line}", 0.7, 2.85 + i * 0.54, 5.3, 0.48, 15, color=MGRAY)

# Right column — what Genesis does
r2 = rect(s, 6.55, 2.1, 6.3, 4.8, fill=INK, line=None)
tb(s, "Genesis Brain Engine", 6.75, 2.25, 5.9, 0.5, 16, bold=True, color=YELLOW)
for i, line in enumerate([
    "Conducts structured brand discovery",
    "Reasons with 25+ yrs of strategy knowledge",
    "Produces three directions simultaneously",
    "Briefs and dispatches specialist agents",
    "Inspects output — iterates until quality met",
    "Retains brand memory across every session",
    "Stays as your creative director — subscription",
]):
    tb(s, f"✓  {line}", 6.75, 2.85 + i * 0.54, 5.9, 0.48, 15, color=INKTEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — How It Works
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "How It Works")

steps = [
    ("01", "Discovery",      "Client answers an AI-guided questionnaire.\nThe Brain Engine conducts intelligent dialogue\nto extract deep brand signals."),
    ("02", "Strategy",       "The Brain Engine reasons against 25+ years\nof brand knowledge and derives three\ncreative directions: anchored, evolved, disruptive."),
    ("03", "Generation",     "Specialist agents — Copy, Visual, Logo,\nPlaybook, Brand Book — are briefed and\norchestrated concurrently."),
    ("04", "Quality Loop",   "The Brain Engine inspects every agent's\noutput against quality rubrics. It iterates\nuntil the threshold is met — not just once."),
    ("05", "Delivery",       "Complete brand delivered: identity, copy,\nbrand book, all assets. Client enters\nongoing consultant mode."),
]
for i, (num, title, body) in enumerate(steps):
    x = 0.4 + i * 2.59
    r = rect(s, x, 1.4, 2.38, 5.4, fill=LGRAY, line=BORDER)
    # number badge
    nb = rect(s, x + 0.85, 1.55, 0.65, 0.5, fill=YELLOW, line=None)
    tb(s, num, x + 0.85, 1.56, 0.65, 0.5, 14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    tb(s, title, x + 0.15, 2.2, 2.1, 0.5, 16, bold=True, color=BLACK)
    tb(s, body, x + 0.15, 2.82, 2.1, 3.5, 13, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Technology: Alpha Architecture (PLACEHOLDER)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Technology")
tb(s, "Alpha Architecture — where it began (~5 years ago)",
   0.55, 1.3, 11.0, 0.5, 18, color=MGRAY)
r = rect(s, 0.5, 1.9, 12.3, 4.5, fill=LGRAY, line=BORDER)
tb(s, "[ INSERT: Alpha Architecture Diagram ]",
   0.5, 3.8, 12.3, 0.8, 22, color=BORDER, align=PP_ALIGN.CENTER, italic=True)
tb(s, "Decision-tree pipeline · validated core proposition · real customer testimonials · proved the market",
   0.5, 6.55, 12.3, 0.45, 13, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Technology: Genesis AI Architecture (PLACEHOLDER)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Technology")
tb(s, "Genesis AI — Production Architecture",
   0.55, 1.3, 11.0, 0.5, 18, color=MGRAY)
r = rect(s, 0.5, 1.9, 12.3, 4.5, fill=LGRAY, line=BORDER)
tb(s, "[ INSERT: Genesis AI Architecture Diagram ]",
   0.5, 3.8, 12.3, 0.8, 22, color=BORDER, align=PP_ALIGN.CENTER, italic=True)
tb(s, "Brain Engine (Claude Opus 4 + extended thinking) · five specialist agents · two-layer knowledge base · Azure cloud infrastructure",
   0.5, 6.55, 12.3, 0.45, 13, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Technology: Brain Engine Deep Dive (PLACEHOLDER)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "The Brain Engine")
tb(s, "Two-layer knowledge architecture — it learns, it reasons, it remembers.",
   0.55, 1.3, 11.0, 0.5, 18, color=MGRAY)
r = rect(s, 0.5, 1.9, 12.3, 4.5, fill=LGRAY, line=BORDER)
tb(s, "[ INSERT: Brain Engine Detail Diagram ]",
   0.5, 3.8, 12.3, 0.8, 22, color=BORDER, align=PP_ALIGN.CENTER, italic=True)
tb(s, "Layer 1: structured YAML reasoning modules distilled from 25+ years of brand strategy  ·  Layer 2: vector store of all past engagements — every brand made sharpens the next",
   0.5, 6.55, 12.3, 0.45, 13, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Market
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Market")

markets = [
    ("5.5M",  "Total Addressable Market",  "New US business registrations per year\n(US Census Bureau, 2023)"),
    ("1.5M",  "Serviceable Market",         "US tech, creative & service startups\nrequiring brand identity"),
    ("150K",  "Obtainable Market",           "10% of serviceable market —\nfastest-growing segment"),
]
for i, (num, label, sub) in enumerate(markets):
    y = 1.5 + i * 1.75
    r = rect(s, 0.5, y, 8.5, 1.55, fill=LGRAY, line=BORDER)
    accent_bar(s, 0.5, y, 1.55)
    tb(s, num,   0.75, y + 0.08, 2.5, 0.8, 44, bold=True, color=BLACK)
    tb(s, label, 3.1,  y + 0.08, 5.7, 0.45, 17, bold=True, color=BLACK)
    tb(s, sub,   3.1,  y + 0.62, 5.7, 0.8,  14, color=DGRAY)

# Additional context right panel
r2 = rect(s, 9.3, 1.5, 3.7, 5.3, fill=INK, line=None)
tb(s, "Why now?", 9.5, 1.65, 3.3, 0.45, 15, bold=True, color=YELLOW)
for i, line in enumerate([
    "AI now capable of genuine\nbrand reasoning — not just\ntemplate matching",
    "",
    "The idea was ahead of the\ntech five years ago.",
    "",
    "The tech has caught up.",
    "",
    "Nobody has built the\nfull loop yet.",
]):
    tb(s, line, 9.5, 2.2 + i * 0.42, 3.3, 0.5, 13, color=INKTEXT if line else INKTEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Business Model
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Business Model")
tb(s, "Four revenue layers — entry, retention, expansion, and growth.",
   0.55, 1.3, 12.0, 0.5, 18, color=DGRAY)

layers = [
    ("Entry",       "Brand Packages",       "$229 / $349 / $849",
     "One-time purchase. Client gets complete brand identity:\nlogo, copy, colour, typography, brand book, all digital assets."),
    ("Retention",   "Subscription",         "$39 – $149 / month",
     "Ongoing access to the AI consultant who knows their brand.\nIncludes asset regeneration, social media graphics, carousels,\npost templates. Higher tiers: social calendar management."),
    ("Expansion",   "Physical Products",    "20–40% margin",
     "Branded business cards, banners, roller banners, t-shirts,\nmugs, and stationery — ordered in-platform, fulfilled via\nprint-on-demand partners. Zero inventory risk."),
    ("Growth",      "AI Social Management", "Premium tier — future",
     "The AI manages the brand's social media presence:\ngenerates, schedules, and publishes content autonomously.\nBrand voice and visual identity always consistent."),
]
for i, (tag, title, price, body) in enumerate(layers):
    y = 1.95 + i * 1.32
    is_future = "future" in price
    bg = RGBColor(0xFF, 0xF9, 0xCC) if not is_future else LGRAY
    r = rect(s, 0.5, y, 12.3, 1.2, fill=bg, line=BORDER)
    accent_bar(s, 0.5, y, 1.2, color=YELLOW if not is_future else BORDER)
    tb(s, tag,   0.75, y + 0.05, 1.5, 0.4, 11, bold=True, color=MGRAY)
    tb(s, title, 0.75, y + 0.45, 2.5, 0.55, 16, bold=True, color=BLACK)
    tb(s, price, 3.5,  y + 0.35, 2.5, 0.55, 18, bold=True,
       color=BLACK if not is_future else MGRAY)
    tb(s, body,  6.3,  y + 0.08, 6.4, 1.05, 13, color=DGRAY if not is_future else MGRAY,
       italic=is_future)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11 — Pricing Detail
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Pricing")

# Brand packages (top half)
tb(s, "Brand Packages — one-time", 0.55, 1.3, 12.0, 0.45, 16, bold=True, color=BLACK)
tiers = [
    ("Essential", "$229",
     ["Full digital brand kit", "Logo (SVG + PNG)", "Brand book (PDF)",
      "Colour & typography", "Copywriting matrix", "Social media graphics",
      "Stationery downloads"]),
    ("Premium", "$349",
     ["Everything in Essential", "Hard copy brand book", "Priority processing", ""]),
    ("Ultimate", "$849",
     ["Everything in Premium", "Full printed stationery pack",
      "Roller banner", "White-glove onboarding"]),
]
for i, (name, price, features) in enumerate(tiers):
    x = 0.5 + i * 4.3
    is_mid = i == 1
    r = rect(s, x, 1.85, 4.0, 2.5, fill=INK if is_mid else LGRAY,
             line=None if is_mid else BORDER, lw=2.0 if is_mid else 1.0)
    tb(s, name,  x + 0.2, 2.0,  3.6, 0.45, 16, bold=True,
       color=YELLOW if is_mid else BLACK)
    tb(s, price, x + 0.2, 2.45, 3.6, 0.6,  28, bold=True,
       color=INKTEXT if is_mid else BLACK)
    for j, feat in enumerate([f for f in features if f][:5]):
        tb(s, f"• {feat}", x + 0.2, 3.1 + j * 0.35, 3.6, 0.32, 12,
           color=INKTEXT if is_mid else DGRAY)

# Subscription tiers (bottom half)
tb(s, "Subscription — monthly", 0.55, 4.55, 12.0, 0.4, 16, bold=True, color=BLACK)
subs = [
    ("Starter",  "$39/mo", "Consultant access · 10 asset regenerations/month"),
    ("Growth",   "$79/mo", "Unlimited consultant · unlimited assets · social media templates & carousels"),
    ("Pro",      "$149/mo","Everything in Growth · AI-managed social media calendar · priority support"),
]
for i, (name, price, desc) in enumerate(subs):
    x = 0.5 + i * 4.3
    r = rect(s, x, 5.05, 4.0, 1.85, fill=LGRAY, line=BORDER)
    tb(s, name,  x + 0.2, 5.15, 2.0, 0.4, 15, bold=True, color=BLACK)
    tb(s, price, x + 2.3, 5.15, 1.7, 0.4, 18, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
    tb(s, desc,  x + 0.2, 5.65, 3.6, 1.1, 13, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12 — Competition
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Competition")

cols = [
    ("Tailor Brands / Looka", "Logo-first tools. Template matching, no strategy, no brand memory, no consultant mode. No subscription with real AI value."),
    ("Canva",                 "Design tool for managing existing brands. No brand creation, no strategy, no discovery. Subscription but no intelligence."),
    ("99designs / Fiverr",    "Human freelancers. Expensive ($500–$5,000+), slow (1–4 weeks), inconsistent. No ongoing relationship."),
    ("Traditional Agency",    "$15K–$75K+, weeks to deliver, entirely inaccessible to startups and small businesses."),
    ("Genesis AI",            "The only platform that combines AI brand strategy, three creative directions, a quality inspection loop, brand memory, physical products, and an ongoing AI consultant — at startup-accessible pricing."),
]
for i, (name, desc) in enumerate(cols):
    y = 1.5 + i * 1.0
    is_us = name == "Genesis AI"
    fill = RGBColor(0xFF, 0xF9, 0xCC) if is_us else LGRAY
    r = rect(s, 0.5, y, 12.3, 0.88, fill=fill,
             line=YELLOW if is_us else BORDER, lw=2.0 if is_us else 1.0)
    tb(s, name, 0.7, y + 0.1, 3.2, 0.65, 16,
       bold=True, color=BLACK)
    tb(s, desc, 4.0, y + 0.1, 8.7, 0.65, 15,
       color=DGRAY, italic=not is_us)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 13 — Team
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Team")

for i, (name, role, bio) in enumerate([
    ("Darrell Irwin",
     "Founder & CEO",
     "25+ years as brand strategist and agency founder. Built and ran Cre8ion, a full-service brand agency. Every framework, rubric, and creative methodology powering the Genesis AI Brain Engine is distilled from his body of work. The intelligence in the system is his."),
    ("Saulo Oliveira",
     "Co-Founder & CTO",
     "Enterprise software architect and AI systems engineer with 20+ years experience. Designed the Genesis AI multi-agent orchestration architecture, the Brain Engine, and the two-layer knowledge base. Previously led platform-level AI and cloud-native architecture at global financial institutions."),
]):
    x = 0.5 + i * 6.5
    r = rect(s, x, 1.5, 6.1, 4.5, fill=LGRAY, line=BORDER)
    accent_bar(s, x, 1.5, 4.5)
    tb(s, name, x + 0.25, 1.65, 5.7, 0.55, 26, bold=True, color=BLACK)
    tb(s, role, x + 0.25, 2.25, 5.7, 0.4,  16, bold=True, color=MGRAY)
    tb(s, bio,  x + 0.25, 2.75, 5.7, 3.0,  15, color=DGRAY)

r3 = rect(s, 0.5, 6.2, 12.3, 0.75, fill=INK, line=None)
tb(s, "$300,000+ in time and expertise invested to date.",
   0.5, 6.3, 12.3, 0.55, 17, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 14 — Traction & Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Traction & Roadmap")

milestones = [
    ("~5 years ago",    "Alpha built & validated",
     "Decision-tree pipeline producing PDF brand books. Deployed, user-tested across real clients spanning professional services, real estate, and recruitment. Customer testimonials on record."),
    ("Today — LIVE",    "Brain Engine operational",
     "Genesis AI Brain Engine is live. Claude Opus 4 with extended thinking, two-layer knowledge base, consultant mode active. Already providing brand insights and proving its intelligence. Demo scheduled."),
    ("Next",            "Agent pipeline integration",
     "Connect specialist agents (Copy, Visual, Logo, Playbook, Brand Book). Test and refine the Brain Engine's ability to brief agents, inspect output quality, and iterate to threshold."),
    ("Then",            "User interface & questionnaire",
     "Build the client-facing discovery experience. AI-guided questionnaire feeding directly into the Brain Engine pipeline."),
    ("End of 2026",     "Beta launch",
     "Full pipeline connected, tested, and refined. Beta opens to early subscribers."),
]
for i, (when, title, body) in enumerate(milestones):
    y = 1.45 + i * 1.1
    is_live = "LIVE" in when
    r = rect(s, 0.5, y, 12.3, 0.98,
             fill=INK if is_live else LGRAY,
             line=YELLOW if is_live else BORDER,
             lw=2.0 if is_live else 1.0)
    tc = YELLOW if is_live else MGRAY
    bc = INKTEXT if is_live else BLACK
    dc = INKTEXT if is_live else DGRAY
    tb(s, when,  0.7,  y + 0.08, 1.9, 0.38, 12, bold=True, color=tc)
    tb(s, title, 0.7,  y + 0.5,  1.9, 0.38, 14, bold=True, color=bc)
    tb(s, body,  2.85, y + 0.08, 9.8, 0.82, 13, color=dc)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 15 — Financial Projections
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Financial Projections")
tb(s, "Three revenue streams · $349 avg brand package · $59/mo avg subscription · 30% physical product margin",
   0.55, 1.25, 12.5, 0.4, 13, color=MGRAY, italic=True)

# Headers
hdrs = ["Year", "Scenario", "Brands/mo", "Packages", "+ Subscriptions", "+ Physical", "Total Revenue"]
col_x = [0.5, 1.3, 2.5, 3.8, 5.5, 7.4, 9.2]
col_w = [0.75, 1.1, 1.2, 1.6, 1.8, 1.7, 3.9]

rh = rect(s, 0.5, 1.75, 12.3, 0.45, fill=INK, line=None)
for ci, (hdr, cx, cw) in enumerate(zip(hdrs, col_x, col_w)):
    tb(s, hdr, cx, 1.77, cw, 0.38, 12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

years_data = [
    ("Y1", [
        ("High", "300",  "$1,258,800", "+$382,000",  "+$72,000",  "$1,712,800"),
        ("Mid",  "150",  "$629,400",   "+$191,000",  "+$36,000",  "$856,400"),
        ("Low",  "75",   "$314,700",   "+$95,000",   "+$18,000",  "$427,700"),
    ]),
    ("Y2", [
        ("High", "1,000","$4,188,000", "+$2,980,000","+$180,000", "$7,348,000"),
        ("Mid",  "500",  "$2,094,000", "+$1,490,000","+$90,000",  "$3,674,000"),
        ("Low",  "250",  "$1,047,000", "+$745,000",  "+$45,000",  "$1,837,000"),
    ]),
    ("Y3", [
        ("High", "5,000","$20,940,000","+$8,530,000","+$500,000", "$29,970,000"),
        ("Mid",  "2,500","$10,470,000","+$4,265,000","+$250,000", "$14,985,000"),
        ("Low",  "1,250","$5,235,000", "+$2,132,000","+$125,000", "$7,492,000"),
    ]),
]

row_y = 2.3
for (yr_label, rows) in years_data:
    h = len(rows) * 0.62
    yb = rect(s, 0.5, row_y, 0.75, h, fill=YELLOW, line=None)
    tb(s, yr_label, 0.5, row_y + h/2 - 0.22, 0.75, 0.44, 16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    for ri, (scen, brands, pkgs, subs, phys, total) in enumerate(rows):
        ry = row_y + ri * 0.62
        bg = LGRAY if ri % 2 == 0 else BG
        rr = rect(s, 1.3, ry, 11.5, 0.58, fill=bg, line=BORDER, lw=0.5)
        vals = [scen, brands, pkgs, subs, phys, total]
        for ci, (val, cx, cw) in enumerate(zip(vals, col_x[1:], col_w[1:])):
            bold = ci == 5
            col = BLACK if not bold else RGBColor(0x11, 0x77, 0x11)
            tb(s, val, cx, ry + 0.1, cw, 0.38, 12, bold=bold, color=col, align=PP_ALIGN.CENTER)
    row_y += h + 0.12

tb(s, "* Investment requirement ~$530K · cost base grows 20% per year · subscriptions assume 25% conversion of brand buyers",
   0.5, 7.1, 12.3, 0.35, 11, color=MGRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 16 — Impact
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_heading(s, "Impact")
tb(s, "What Genesis AI means for American entrepreneurs.",
   0.55, 1.3, 11.5, 0.5, 20, color=DGRAY, italic=True)

points = [
    ("Democratizing Entrepreneurship",
     "5.5M Americans start a business every year. Fewer than 10% can afford professional branding. Genesis AI removes the most-cited barrier to early-stage success — permanently."),
    ("Transforming the Economics",
     "What costs $15,000–$75,000 at an agency and takes weeks now takes minutes and costs hundreds. The value-to-price gap is not incremental — it is categorical."),
    ("AI That Gets Smarter Over Time",
     "Every engagement strengthens Layer 2 — the vector store of brand knowledge. The system that serves the 10,000th client is meaningfully smarter than the one that served the first."),
    ("A New Creative Economy",
     "Creatives contribute knowledge and earn when it's applied. A responsible AI model that amplifies the creative industry rather than replacing it."),
]
for i, (title, body) in enumerate(points):
    x = 0.5 + (i % 2) * 6.45
    y = 2.0 + (i // 2) * 2.5
    r = rect(s, x, y, 6.1, 2.2, fill=LGRAY, line=BORDER)
    accent_bar(s, x, y, 2.2)
    tb(s, title, x + 0.22, y + 0.1, 5.7, 0.5,  17, bold=True, color=BLACK)
    tb(s, body,  x + 0.22, y + 0.72, 5.7, 1.35, 14, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 17 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0.0, 0.0, 0.18, 7.5, fill=YELLOW, line=None)
tb(s, "Creating the world's most intelligent\nbrand agency will inspire people\nto design their future.",
   0.55, 1.3, 11.5, 2.8, 36, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
tb(s, "The automation race has just begun.",
   0.55, 4.4, 11.5, 0.75, 28, color=MGRAY, align=PP_ALIGN.LEFT)
tb(s, "Brain Engine: live  ·  Demo: scheduled  ·  Beta: end of 2026",
   0.55, 5.4, 11.5, 0.5,  18, bold=True, color=BLACK)
tb(s, "genesisbrands.ai",
   0.55, 6.4, 5.0,  0.5,  15, color=MGRAY)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/saulo/genesis-ai/docs/eb2-niw/genesis-ai-presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
