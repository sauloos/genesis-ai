"""Extract text from PowerPoint files.

Text-based slides: read directly via python-pptx.
Image-based slides: send each slide image to GPT-4o vision for OCR.
"""

import base64
import os
from pathlib import Path


def extract(file_path: str) -> dict:
    """
    Extract all text from a .pptx file.
    Falls back to vision OCR for slides that contain no extractable text.
    Returns dict with keys: text, title, author.
    """
    from pptx import Presentation
    from openai import OpenAI

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    path = Path(file_path)
    prs = Presentation(str(path))

    core = prs.core_properties
    author = core.author or ""

    slide_texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        text = _extract_text_shapes(slide)
        if text:
            slide_texts.append(text)
        else:
            # Slide is image-based — OCR with vision
            images = _extract_images(slide)
            if images:
                print(f"  Slide {slide_num}: OCR via vision ({len(images)} image(s))...")
                ocr_text = _ocr_images(client, images, slide_num)
                if ocr_text:
                    slide_texts.append(ocr_text)

    text = "\n\n".join(slide_texts)
    title = path.stem
    return {"text": text, "title": title, "author": author}


def _extract_text_shapes(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
    return " ".join(parts)


def _extract_images(slide) -> list[bytes]:
    """Return raw image bytes for every picture shape on the slide."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                images.append(shape.image.blob)
            except Exception:
                pass
    return images


def _ocr_images(client, image_blobs: list[bytes], slide_num: int) -> str:
    """Send images to GPT-4o vision and return extracted text."""
    content = [
        {
            "type": "text",
            "text": (
                "Extract all text visible in this presentation slide. "
                "Return only the text content, preserving structure with newlines. "
                "Do not describe the layout or design — only the words."
            ),
        }
    ]
    for blob in image_blobs:
        b64 = base64.b64encode(blob).decode()
        # Detect format from magic bytes
        mime = "image/jpeg"
        if blob[:4] == b'\x89PNG':
            mime = "image/png"
        elif blob[:4] == b'GIF8':
            mime = "image/gif"
        elif blob[:4] == b'RIFF':
            mime = "image/webp"
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"},
        })

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": content}],
        max_tokens=2000,
    )
    return response.choices[0].message.content.strip()
